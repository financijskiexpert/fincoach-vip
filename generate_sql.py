import sys
import re
sys.stdout.reconfigure(encoding='utf-8')

from pptx import Presentation

files = [
    (r'C:\Users\Lenovo\Downloads\Put prema financijskoj slobodi - DELAVNICA\Financijske strategije za osvajanje ciljeva 0-30\Financijske strategije za osvajanje ciljeva0-30.pptx', 1, 30),
    (r'C:\Users\Lenovo\Downloads\Put prema financijskoj slobodi - DELAVNICA\Financijske strategije za osvajanje ciljeva 30-60\Financijske strategije za osvajanje ciljeva30-60.pptx', 31, 60),
    (r'C:\Users\Lenovo\Downloads\Put prema financijskoj slobodi - DELAVNICA\Financijske strategije za osvajanje ciljeva 60-90\Financijske strategije za osvajanje ciljeva60-90.pptx', 61, 90),
]

sections = {
    range(1, 31): 'Faza 1: Put prema financijskoj slobodi',
    range(31, 61): 'Faza 2: Financijske strategije za osvajanje ciljeva',
    range(61, 91): 'Faza 3: Financijska neovisnost',
}

def get_section(day):
    for r, s in sections.items():
        if day in r:
            return s
    return 'Ostalo'

lessons = {}

for path, day_start, day_end in files:
    prs = Presentation(path)
    for slide in prs.slides:
        all_text = []
        for shape in slide.shapes:
            if hasattr(shape, 'text') and shape.text.strip():
                t = shape.text.strip()
                if t != 'Volim Svoj Novac':
                    all_text.append(t)

        full_text = ' '.join(all_text)
        day_match = re.search(r'(\d+)\.\s*Dan', full_text)
        if not day_match:
            continue

        day_num = int(day_match.group(1))
        if day_num < day_start or day_num > day_end:
            continue
        if day_num in lessons:
            continue

        content = re.sub(r'\d+\.\s*Dan', '', full_text).strip()
        first_sentence = re.split(r'[.!?]', content)[0].strip()
        # Clean up quotes and special chars
        first_sentence = first_sentence.strip('"').strip("'").strip()
        if len(first_sentence) > 100:
            first_sentence = first_sentence[:97] + '...'
        if len(first_sentence) < 10:
            first_sentence = f'Dan {day_num}'

        lessons[day_num] = first_sentence

course_id = 'f1fb4513-11e4-4557-b771-56a73e1ec03e'

# Generate SQL
lines = []
lines.append('-- Brisanje starih lekcija i ubacivanje novih s pravim naslovima')
lines.append(f"DELETE FROM lessons WHERE course_id = '{course_id}';")
lines.append('')
lines.append('INSERT INTO lessons (course_id, day_number, title, section, sort_order, duration_seconds, is_published) VALUES')

rows = []
for day in sorted(lessons.keys()):
    title = lessons[day].replace("'", "''")
    section = get_section(day).replace("'", "''")
    rows.append(f"('{course_id}', {day}, '{title}', '{section}', {day}, 0, true)")

lines.append(',\n'.join(rows) + ';')

sql = '\n'.join(lines)

with open(r'E:\fincoach-vip\insert_lessons_final.sql', 'w', encoding='utf-8') as f:
    f.write(sql)

print(f'SQL generiran za {len(lessons)} lekcij.')
print('Datoteka: E:\\fincoach-vip\\insert_lessons_final.sql')
