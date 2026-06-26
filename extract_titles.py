from pptx import Presentation

files = [
    (r'C:\Users\Lenovo\Downloads\Put prema financijskoj slobodi - DELAVNICA\Financijske strategije za osvajanje ciljeva 0-30\Financijske strategije za osvajanje ciljeva0-30.pptx', '0-30'),
    (r'C:\Users\Lenovo\Downloads\Put prema financijskoj slobodi - DELAVNICA\Financijske strategije za osvajanje ciljeva 30-60\Financijske strategije za osvajanje ciljeva30-60.pptx', '30-60'),
    (r'C:\Users\Lenovo\Downloads\Put prema financijskoj slobodi - DELAVNICA\Financijske strategije za osvajanje ciljeva 60-90\Financijske strategije za osvajanje ciljeva60-90.pptx', '60-90'),
]

for path, label in files:
    print(f'\n=== PPTX {label} ===')
    prs = Presentation(path)
    for i, slide in enumerate(prs.slides, 1):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, 'text') and shape.text.strip():
                t = shape.text.strip().replace('\n', ' | ')
                if t != 'Volim Svoj Novac':
                    texts.append(t[:150])
        summary = ' // '.join(texts[:4]) if texts else '[prazno]'
        print(f'  Slide {i}: {summary[:200]}')
