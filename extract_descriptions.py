"""Extract descriptions from PPTX files and output as JSON."""
import sys
import re
import json
sys.stdout.reconfigure(encoding='utf-8')

from pptx import Presentation

files = [
    (r'C:\Users\Lenovo\Downloads\Put prema financijskoj slobodi - DELAVNICA\Financijske strategije za osvajanje ciljeva 0-30\Financijske strategije za osvajanje ciljeva0-30.pptx', 1, 30),
    (r'C:\Users\Lenovo\Downloads\Put prema financijskoj slobodi - DELAVNICA\Financijske strategije za osvajanje ciljeva 30-60\Financijske strategije za osvajanje ciljeva30-60.pptx', 31, 60),
    (r'C:\Users\Lenovo\Downloads\Put prema financijskoj slobodi - DELAVNICA\Financijske strategije za osvajanje ciljeva 60-90\Financijske strategije za osvajanje ciljeva60-90.pptx', 61, 90),
]

descriptions = {}

for path, day_start, day_end in files:
    try:
        prs = Presentation(path)
    except Exception as e:
        print(f"ERROR opening {path}: {e}", file=sys.stderr)
        continue

    for slide in prs.slides:
        all_text = []
        for shape in slide.shapes:
            if hasattr(shape, 'text') and shape.text.strip():
                t = shape.text.strip()
                if t != 'Volim Svoj Novac':
                    all_text.append(t)

        full_text = '\n'.join(all_text)
        day_match = re.search(r'(\d+)\.\s*Dan', full_text)
        if not day_match:
            continue
        day_num = int(day_match.group(1))
        if day_num < day_start or day_num > day_end:
            continue
        if day_num in descriptions:
            continue

        # Remove day marker line and keep the rest as description
        lines = [l.strip() for l in full_text.split('\n') if l.strip()]
        # Skip lines that are just the day marker
        cleaned = []
        for l in lines:
            if re.match(r'^\d+\.\s*Dan\s*$', l):
                continue
            cleaned.append(l)
        desc = '\n\n'.join(cleaned)
        # Limit to 1500 chars
        if len(desc) > 1500:
            desc = desc[:1497] + '...'
        descriptions[day_num] = desc

print(json.dumps(descriptions, ensure_ascii=False, indent=2))
