import sys
import re
sys.stdout.reconfigure(encoding='utf-8')

from pptx import Presentation

files = [
    (r'C:\Users\Lenovo\Downloads\Put prema financijskoj slobodi - DELAVNICA\Financijske strategije za osvajanje ciljeva 0-30\Financijske strategije za osvajanje ciljeva0-30.pptx', 1, 30),
    (r'C:\Users\Lenovo\Downloads\Put prema financijskoj slobodi - DELAVNICA\Financijske strategije za osvajanje ciljeva 30-60\Financijske strategije za osvajanje ciljeva30-60.pptx', 31, 60),
    (r'C:\Users\Lenovo\Downloads\Put prema financijskoj slobodi - DELAVNICA\Financijske strategije za osvajanje ciljeva 60-90\Financijske strategije za osvajanje ciljeva60-90.pptx', 61, 90),
]

lessons = {}

for path, day_start, day_end in files:
    prs = Presentation(path)
    for slide in prs.slides:
        all_text = []
        for shape in slide.shapes:
            if hasattr(shape, 'text') and shape.text.strip():
                t = shape.text.strip()
                if t != 'Volim Svoj Novac':
                    all_text.append(t)

        full_text = ' '.join(all_text)

        # Find day number
        day_match = re.search(r'(\d+)\.\s*Dan', full_text)
        if not day_match:
            continue

        day_num = int(day_match.group(1))
        if day_num < day_start or day_num > day_end:
            continue

        # Only take first occurrence per day
        if day_num in lessons:
            continue

        # Extract meaningful title from first part of text (before the day marker)
        content = re.sub(r'\d+\.\s*Dan', '', full_text).strip()
        # Get first sentence
        first_sentence = re.split(r'[.!?]', content)[0].strip()
        if len(first_sentence) > 80:
            first_sentence = first_sentence[:77] + '...'
        if len(first_sentence) < 10:
            first_sentence = content[:80].strip()

        lessons[day_num] = first_sentence

# Print all 90 lessons
for day in sorted(lessons.keys()):
    print(f'{day}|{lessons[day]}')
